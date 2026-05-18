from __future__ import annotations

import importlib.util
from pathlib import Path

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


HERE = Path(__file__).resolve().parent
SRC = HERE / "draw_data_workflow_reference_v3.py"
ASSETS = HERE / "ppt_editable_assets"
OUT = HERE / "data_workflow_reference_v3_editable.pptx"

spec = importlib.util.spec_from_file_location("workflow_fig", SRC)
workflow = importlib.util.module_from_spec(spec)
spec.loader.exec_module(workflow)

COLORS = {
    "ink": "1F2933",
    "blue": "315F8C",
    "teal": "4D8E88",
    "orange": "B56B1D",
    "red": "A84A4A",
    "green": "4F7D55",
    "gray": "7F8A93",
    "bg": "FCFDFE",
    "white": "FFFFFF",
}


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def save_asset(name: str, draw_func, size=(2.0, 1.0), dpi=300) -> Path:
    ASSETS.mkdir(parents=True, exist_ok=True)
    path = ASSETS / f"{name}.png"
    fig = plt.figure(figsize=size, dpi=dpi)
    ax = fig.add_axes([0, 0, 1, 1])
    draw_func(ax)
    fig.savefig(path, dpi=dpi, transparent=False, bbox_inches="tight", pad_inches=0.01)
    plt.close(fig)
    return path


def set_text(shape, value, size=7, bold=False, color="ink"):
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = value
    run.font.name = "SimSun"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(COLORS[color])


def add_box(slide, x, y, w, h, label, line="blue", fill="white", size=7, bold=False, width=0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS[fill])
    shape.line.color.rgb = rgb(COLORS[line])
    shape.line.width = Pt(width)
    set_text(shape, label, size=size, bold=bold)
    return shape


def add_outer(slide, x, y, w, h, title, line):
    shape = add_box(slide, x, y, w, h, "", line=line, fill="bg", width=1.0)
    shape.line.dash_style = 4
    title_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.035), Inches(w), Inches(0.22))
    set_text(title_box, title, size=10, bold=False)
    return shape


def add_text(slide, x, y, w, h, label, size=6, bold=False, color="ink"):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(shape, label, size=size, bold=bold, color=color)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color="gray", width=1.0):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    connector.line.color.rgb = rgb(COLORS[color])
    connector.line.width = Pt(width)
    connector.line.end_arrowhead = True
    return connector


def add_picture(slide, path, x, y, w, h):
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def build_assets():
    boundary, holes, odi, params, stats, coal = workflow.load_data()
    return {
        "boundary": save_asset("boundary", lambda ax: workflow.map_ax(ax, boundary, holes, "边界/钻孔"), size=(1.8, 0.95)),
        "lith": save_asset("lithology", lambda ax: workflow.lith_ax(ax), size=(1.8, 0.95)),
        "param": save_asset("param_field", lambda ax: workflow.field_ax(ax, params, boundary, "Hi", "参数场", "YlGnBu"), size=(1.8, 0.95)),
        "odi": save_asset("odi_field", lambda ax: workflow.field_ax(ax, odi, boundary, "ODI_norm", "ODI场", "YlOrRd"), size=(1.8, 0.95)),
        "spatial": save_asset(
            "spatial_clip",
            lambda ax: workflow.map_ax(ax, boundary, holes, "坐标/边界", odi["ODI_norm"], odi),
            size=(2.0, 1.10),
        ),
        "hi": save_asset("hi_interp", lambda ax: workflow.field_ax(ax, params, boundary, "Hi", "Hi插值场", "YlGnBu"), size=(2.0, 1.10)),
        "odi_norm": save_asset("odi_norm", lambda ax: workflow.field_ax(ax, odi, boundary, "ODI_norm", "ODI场", "YlOrRd"), size=(2.0, 1.10)),
        "valid": save_asset("valid_domain", lambda ax: workflow.map_ax(ax, boundary, holes, "有效域"), size=(1.8, 0.95)),
        "coal": save_asset("coal_resource", lambda ax: workflow.field_ax(ax, coal, boundary, "coal_thickness", "煤厚资源", "YlOrBr"), size=(1.8, 0.95)),
        "odi_layer": save_asset("odi_layer", lambda ax: workflow.field_ax(ax, odi, boundary, "ODI_norm", "ODI扰动", "YlOrRd"), size=(1.8, 0.95)),
        "candidate": save_asset("candidate_layer", lambda ax: workflow.candidate_ax(ax, boundary), size=(1.8, 0.95)),
        "stats": save_asset("stats", lambda ax: workflow.stats_ax(ax, stats), size=(2.3, 1.05)),
    }, len(holes), len(odi)


def build_pptx():
    assets, n_holes, n_odi = build_assets()

    prs = Presentation()
    prs.slide_width = Inches(7.2)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(COLORS["white"])

    add_outer(slide, 0.25, 0.27, 6.70, 1.18, "数据来源与分类", "blue")
    add_outer(slide, 0.25, 1.75, 6.70, 1.78, "空间处理与指标构建", "teal")
    add_outer(slide, 0.25, 3.85, 6.70, 1.36, "图层输出与模型输入", "orange")
    add_outer(slide, 0.25, 5.55, 6.70, 1.50, "规划应用与决策输出", "red")
    add_arrow(slide, 3.60, 1.45, 3.60, 1.72, "gray", 1.0)
    add_arrow(slide, 3.60, 3.53, 3.60, 3.82, "gray", 1.0)
    add_arrow(slide, 3.60, 5.21, 3.60, 5.52, "gray", 1.0)

    # Data source layer.
    source_boxes = [
        (0.43, "矿井地质资料"),
        (1.80, "钻孔资料"),
        (3.17, "采区设计图件"),
        (4.54, "扰动评价结果"),
    ]
    for x, label in source_boxes:
        add_box(slide, x, 0.62, 1.16, 0.22, label, "blue", size=6.3, width=0.9)
    add_picture(slide, assets["boundary"], 0.58, 0.89, 0.86, 0.49)
    add_picture(slide, assets["lith"], 1.90, 0.89, 0.92, 0.49)
    add_picture(slide, assets["param"], 3.34, 0.89, 0.80, 0.49)
    add_picture(slide, assets["odi"], 4.65, 0.89, 0.80, 0.49)
    add_arrow(slide, 5.50, 1.12, 5.75, 1.12, "blue", 1.0)
    for x, y, label in [
        (5.78, 0.86, "基础地质"),
        (6.40, 0.86, "工程约束"),
        (5.78, 1.14, "覆岩扰动"),
        (6.40, 1.14, "规划参数"),
    ]:
        add_box(slide, x, y, 0.52, 0.20, label, "green", size=5.0, width=0.65)
    add_text(slide, 5.82, 1.34, 1.05, 0.10, f"钻孔 {n_holes} 个  |  ODI点 {n_odi} 个  |  参数 4 类", size=4.2, color="ink")

    # Spatial processing layer.
    modules = [
        (0.47, "坐标统一与裁剪", assets["spatial"], ("格式转换", "边界裁剪")),
        (2.53, "参数场插值", assets["hi"], ("插值计算", "网格构建")),
        (4.59, "ODI归一化", assets["odi_norm"], ("指标归一化", "图层叠加")),
    ]
    for x, title, img, ops in modules:
        frame = add_box(slide, x, 2.10, 1.63, 1.18, "", "teal", fill="white", width=0.65)
        frame.line.dash_style = 3
        add_text(slide, x, 2.18, 1.63, 0.18, title, size=7.2)
        add_picture(slide, img, x + 0.30, 2.42, 1.00, 0.48)
        add_box(slide, x + 0.27, 3.00, 0.48, 0.20, ops[0], "teal", size=5.2, width=0.6)
        add_box(slide, x + 0.88, 3.00, 0.48, 0.20, ops[1], "teal", size=5.2, width=0.6)
    add_arrow(slide, 2.22, 2.69, 2.46, 2.69, "teal", 1.0)
    add_arrow(slide, 4.28, 2.69, 4.52, 2.69, "teal", 1.0)

    # Layer output.
    outputs = [
        (0.45, assets["valid"], "有效布置域图层"),
        (1.95, assets["coal"], "煤厚资源图层"),
        (3.45, assets["odi_layer"], "ODI扰动图层"),
        (4.95, assets["candidate"], "候选方案图层"),
    ]
    for x, img, label in outputs:
        add_picture(slide, img, x, 4.30, 0.82, 0.43)
        add_box(slide, x - 0.10, 4.78, 1.02, 0.23, label, "orange", size=5.5, bold=label.startswith("ODI"), width=0.75)
    add_arrow(slide, 6.08, 4.53, 6.35, 4.53, "orange", 1.0)
    add_box(slide, 6.38, 4.35, 0.42, 0.50, "模型\n输入", "orange", size=6.0, bold=True, width=0.75)

    # Planning layer.
    steps = ["方案生成", "约束筛选", "指标统计", "方案比选"]
    for i, step in enumerate(steps):
        x = 0.45 + i * 0.96
        add_box(slide, x, 6.25, 0.70, 0.30, step, "red", size=6.6, bold=i in (0, 3), width=0.75)
        if i < 3:
            add_arrow(slide, x + 0.73, 6.40, x + 0.91, 6.40, "red", 1.0)
    add_picture(slide, assets["stats"], 4.23, 6.06, 1.35, 0.62)
    add_arrow(slide, 5.68, 6.33, 5.94, 6.33, "red", 1.0)
    decisions = [
        ("A  效率优先", "DDEBFA"),
        ("B  资源优先", "E5F0D9"),
        ("C  低扰动", "FFF0D5"),
        ("D  不推荐", "F2DCDC"),
    ]
    for i, (label, fill) in enumerate(decisions):
        box_shape = add_box(slide, 6.02, 5.88 + i * 0.25, 0.70, 0.18, label, "red", fill="white", size=4.8, bold=i == 1, width=0.50)
        box_shape.fill.solid()
        box_shape.fill.fore_color.rgb = rgb(fill)
        box_shape.line.dash_style = 3

    add_text(slide, 0.90, 7.22, 5.40, 0.12, "注：缩略图由采区边界、钻孔坐标、钻孔分层、ODI评价点与方案统计数据生成；比例尺仅用于流程说明。", size=4.8, color="gray")

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print(build_pptx())
