from __future__ import annotations

from pathlib import Path

import matplotlib as mpl
import matplotlib.pyplot as plt
from matplotlib.font_manager import FontProperties
from matplotlib.patches import FancyArrowPatch, Rectangle
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


OUT = Path(__file__).resolve().parent
FONT = FontProperties(fname=r"C:\Windows\Fonts\simsun.ttc")
FONT_BOLD = FontProperties(fname=r"C:\Windows\Fonts\simhei.ttf")

mpl.rcParams["font.family"] = "sans-serif"
mpl.rcParams["pdf.fonttype"] = 42
mpl.rcParams["ps.fonttype"] = 42
mpl.rcParams["svg.fonttype"] = "none"
mpl.rcParams["axes.unicode_minus"] = False

INK = "#222222"
GRAY = "#707070"
LIGHT = "#F7F8F9"
GRID = "#C9D0D5"
BLUE = "#2F5F8F"
GREEN = "#4F7F52"
TEAL = "#4F8E8A"
ORANGE = "#B96B21"
RED = "#A74B4B"
ODI = "#A84600"


def hex_to_rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def add_rect(ax, x, y, w, h, label="", edge=GRAY, fill="#FFFFFF", lw=0.8, fs=7.4, bold=False, color=INK):
    ax.add_patch(Rectangle((x, y), w, h, facecolor=fill, edgecolor=edge, linewidth=lw, zorder=2))
    if label:
        ax.text(
            x + w / 2,
            y + h / 2,
            label,
            ha="center",
            va="center",
            fontproperties=FONT_BOLD if bold else FONT,
            fontsize=fs,
            color=color,
            zorder=3,
        )


def add_arrow(ax, start, end, color="#8A9AA3", lw=0.75, scale=7.0):
    ax.add_patch(
        FancyArrowPatch(
            start,
            end,
            arrowstyle="-|>",
            mutation_scale=scale,
            linewidth=lw,
            color=color,
            shrinkA=0,
            shrinkB=0,
            zorder=4,
        )
    )


def draw_matplotlib() -> None:
    fig, ax = plt.subplots(figsize=(7.20, 4.72), dpi=300)
    fig.patch.set_facecolor("white")
    ax.set_xlim(0, 10.8)
    ax.set_ylim(0, 7.08)
    ax.axis("off")

    rows = [
        (5.82, 0.82, "数据来源层", "Data source", BLUE),
        (4.68, 0.82, "数据分类层", "Classification", GREEN),
        (3.26, 1.05, "空间处理层", "Spatial processing", TEAL),
        (2.03, 0.82, "图层输出层", "Layer output", ORANGE),
        (0.88, 0.82, "规划应用层", "Planning", RED),
    ]

    for idx, (y, h, cn, en, accent) in enumerate(rows, start=1):
        add_rect(ax, 0.28, y, 10.22, h, edge=GRID, fill=LIGHT, lw=0.55)
        ax.add_patch(Rectangle((0.28, y), 0.08, h, facecolor=accent, edgecolor=accent, linewidth=0, zorder=3))
        ax.text(0.52, y + h * 0.62, f"{idx}. {cn}", ha="left", va="center", fontproperties=FONT_BOLD, fontsize=7.5, color=INK)
        ax.text(0.52, y + h * 0.32, en, ha="left", va="center", fontproperties=FONT, fontsize=5.5, color="#666666")
        ax.plot([2.0, 2.0], [y + 0.12, y + h - 0.12], color="#D8DEE2", lw=0.55, zorder=2)

    # Source data.
    l1 = ["矿井地质资料", "采区设计图件", "钻孔资料", "覆岩扰动\n评价结果", "工作面规划参数"]
    box_w, gap, x0 = 1.52, 0.18, 2.25
    for i, label in enumerate(l1):
        add_rect(ax, x0 + i * (box_w + gap), 6.04, box_w, 0.38, label, edge=BLUE, lw=0.75, fs=6.55)

    # Classification.
    l2 = ["基础地质数据", "工程约束数据", "覆岩扰动数据", "规划参数数据"]
    box_w, gap, x0 = 1.82, 0.33, 2.43
    for i, label in enumerate(l2):
        add_rect(ax, x0 + i * (box_w + gap), 4.90, box_w, 0.38, label, edge=GREEN, lw=0.75, fs=6.9)

    # Processing chain.
    l3 = ["坐标统一", "格式转换", "边界裁剪", "插值计算", "指标归一化", "图层叠加"]
    box_w, gap, x0, y = 1.18, 0.13, 2.22, 3.66
    for i, label in enumerate(l3):
        add_rect(ax, x0 + i * (box_w + gap), y, box_w, 0.36, label, edge=TEAL, lw=0.75, fs=6.7)
        if i < len(l3) - 1:
            sx = x0 + i * (box_w + gap) + box_w + 0.02
            add_arrow(ax, (sx, y + 0.18), (sx + gap - 0.04, y + 0.18), color=TEAL, lw=0.6, scale=5.0)
    ax.text(2.22, 3.43, "统一空间参考、数据尺度与指标量纲，形成可叠加的空间分析基础", ha="left", va="center", fontproperties=FONT, fontsize=5.65, color="#666666")

    # Output layers.
    l4 = [("有效布置域图层", ORANGE), ("煤厚资源图层", ORANGE), ("ODI扰动图层", ODI), ("候选方案图层", ORANGE)]
    box_w, gap, x0 = 1.82, 0.33, 2.43
    for i, (label, edge) in enumerate(l4):
        add_rect(ax, x0 + i * (box_w + gap), 2.25, box_w, 0.38, label, edge=edge, lw=0.85 if edge == ODI else 0.75, fs=6.9, bold=edge == ODI)

    # Planning applications.
    l5 = ["方案生成", "约束筛选", "指标统计", "方案比选"]
    for i, label in enumerate(l5):
        add_rect(ax, x0 + i * (box_w + gap), 1.10, box_w, 0.38, label, edge=RED, lw=0.75, fs=7.0)

    # Vertical connectors.
    for y1, y2 in [(5.82, 5.50), (4.68, 4.31), (3.26, 2.85), (2.03, 1.70)]:
        add_arrow(ax, (6.24, y1), (6.24, y2), color="#909EA5", lw=0.75, scale=8.0)

    # Subtle many-to-one guide from classified data to the processing chain.
    for x in [3.34, 5.49, 7.64, 9.79]:
        ax.plot([x, x, 6.24], [4.68, 4.48, 4.31], color="#D4DBDF", lw=0.45, zorder=1)

    ax.text(10.42, 0.35, "ODI: overburden disturbance index", ha="right", va="center", fontproperties=FONT, fontsize=5.4, color="#666666")

    OUT.mkdir(parents=True, exist_ok=True)
    stem = OUT / "data_workflow_ppt_scientific"
    fig.savefig(stem.with_suffix(".svg"), bbox_inches="tight", pad_inches=0.035)
    fig.savefig(stem.with_suffix(".pdf"), bbox_inches="tight", pad_inches=0.035)
    fig.savefig(stem.with_suffix(".png"), dpi=900, bbox_inches="tight", pad_inches=0.035)
    plt.close(fig)


def ppt_add_text(shape, text: str, font_size: float, bold: bool = False, color: str = INK):
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "SimSun"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(color)


def ppt_box(slide, x, y, w, h, text, edge=GRAY, fill="#FFFFFF", fs=9, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill)
    shape.line.color.rgb = hex_to_rgb(edge)
    shape.line.width = Pt(0.8)
    ppt_add_text(shape, text, fs, bold, INK)
    return shape


def ppt_line(slide, x1, y1, x2, y2, color="#8A9AA3"):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = hex_to_rgb(color)
    connector.line.width = Pt(0.8)
    connector.line.end_arrowhead = True
    return connector


def draw_pptx() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb("#FFFFFF")

    rows = [
        (0.62, 1.00, "1. 数据来源层", "Data source", BLUE),
        (1.86, 1.00, "2. 数据分类层", "Classification", GREEN),
        (3.10, 1.22, "3. 空间处理层", "Spatial processing", TEAL),
        (4.55, 1.00, "4. 图层输出层", "Layer output", ORANGE),
        (5.80, 1.00, "5. 规划应用层", "Planning", RED),
    ]
    for y, h, cn, en, accent in rows:
        ppt_box(slide, 0.42, y, 12.52, h, "", edge="#C9D0D5", fill="#F7F8F9", fs=1)
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(y), Inches(0.09), Inches(h))
        strip.fill.solid()
        strip.fill.fore_color.rgb = hex_to_rgb(accent)
        strip.line.color.rgb = hex_to_rgb(accent)
        label = slide.shapes.add_textbox(Inches(0.70), Inches(y + 0.18), Inches(1.62), Inches(0.30))
        ppt_add_text(label, cn, 10.5, True, INK)
        sub = slide.shapes.add_textbox(Inches(0.70), Inches(y + 0.54), Inches(1.62), Inches(0.25))
        ppt_add_text(sub, en, 8.0, False, "#666666")
        sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.42), Inches(y + 0.13), Inches(0.01), Inches(h - 0.26))
        sep.fill.solid()
        sep.fill.fore_color.rgb = hex_to_rgb("#D8DEE2")
        sep.line.color.rgb = hex_to_rgb("#D8DEE2")

    # Content boxes.
    for i, label in enumerate(["矿井地质资料", "采区设计图件", "钻孔资料", "覆岩扰动\n评价结果", "工作面规划参数"]):
        ppt_box(slide, 2.72 + i * 1.95, 0.91, 1.70, 0.40, label, edge=BLUE, fs=8.5)

    for i, label in enumerate(["基础地质数据", "工程约束数据", "覆岩扰动数据", "规划参数数据"]):
        ppt_box(slide, 2.94 + i * 2.32, 2.17, 1.95, 0.40, label, edge=GREEN, fs=9.0)

    proc = ["坐标统一", "格式转换", "边界裁剪", "插值计算", "指标归一化", "图层叠加"]
    for i, label in enumerate(proc):
        x = 2.72 + i * 1.50
        ppt_box(slide, x, 3.50, 1.22, 0.38, label, edge=TEAL, fs=8.5)
        if i < len(proc) - 1:
            ppt_line(slide, x + 1.22, 3.69, x + 1.45, 3.69, color=TEAL)
    note = slide.shapes.add_textbox(Inches(2.72), Inches(3.92), Inches(6.80), Inches(0.26))
    ppt_add_text(note, "统一空间参考、数据尺度与指标量纲，形成可叠加的空间分析基础", 8.0, False, "#666666")

    for i, (label, edge, bold) in enumerate([("有效布置域图层", ORANGE, False), ("煤厚资源图层", ORANGE, False), ("ODI扰动图层", ODI, True), ("候选方案图层", ORANGE, False)]):
        ppt_box(slide, 2.94 + i * 2.32, 4.86, 1.95, 0.40, label, edge=edge, fs=9.0, bold=bold)

    for i, label in enumerate(["方案生成", "约束筛选", "指标统计", "方案比选"]):
        ppt_box(slide, 2.94 + i * 2.32, 6.10, 1.95, 0.40, label, edge=RED, fs=9.2, bold=True)

    for y1, y2 in [(1.62, 1.86), (2.86, 3.10), (4.32, 4.55), (5.55, 5.80)]:
        ppt_line(slide, 7.10, y1, 7.10, y2)

    foot = slide.shapes.add_textbox(Inches(9.0), Inches(6.92), Inches(3.72), Inches(0.24))
    ppt_add_text(foot, "ODI: overburden disturbance index", 7.5, False, "#666666")

    prs.save(OUT / "data_workflow_ppt_scientific_editable.pptx")


if __name__ == "__main__":
    draw_matplotlib()
    draw_pptx()
